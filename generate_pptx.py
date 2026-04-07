#!/usr/bin/env python3
"""Physical AI Workshop 프레젠테이션 PPTX 생성기"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG_DARK   = RGBColor(0x06, 0x08, 0x0F)
SURFACE   = RGBColor(0x0C, 0x10, 0x22)
BORDER    = RGBColor(0x1E, 0x29, 0x3B)
TEXT      = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DIM  = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED= RGBColor(0x64, 0x74, 0x8B)
BLUE      = RGBColor(0x60, 0xA5, 0xFA)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
PINK      = RGBColor(0xF4, 0x72, 0xB6)
GREEN     = RGBColor(0x34, 0xD3, 0x99)
ORANGE    = RGBColor(0xFB, 0x92, 0x3C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=TEXT, bold=False, spacing=1.2, font_name='Calibri'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, clr, b) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = b
        p.font.name = 'Calibri'
        p.space_after = Pt(sz * (spacing - 1) + 2)
    return txbox

def add_footer(slide, left_text="Physical AI Workshop", right_text=""):
    add_text(slide, Inches(0.6), Inches(6.9), Inches(4), Inches(0.4),
             left_text, font_size=9, color=TEXT_MUTED)
    add_text(slide, Inches(8.5), Inches(6.9), Inches(4), Inches(0.4),
             right_text, font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

def add_slide_num(slide, num):
    add_text(slide, Inches(12.2), Inches(0.3), Inches(0.8), Inches(0.3),
             f"{num:02d}", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT,
             font_name='Consolas')

def add_tag(slide, left, top, text, color=PURPLE):
    shape = add_shape(slide, left, top, Inches(len(text)*0.12 + 0.6), Inches(0.35),
                      fill_color=RGBColor(color[0]//6, color[1]//6, color[2]//6),
                      border_color=RGBColor(color[0]//3, color[1]//3, color[2]//3))
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(9)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Consolas'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(2)
    return shape

def add_stat_card(slide, left, top, value, label, color=PURPLE, width=Inches(2.4)):
    card = add_shape(slide, left, top, width, Inches(1.1), fill_color=SURFACE, border_color=BORDER)
    add_text(slide, left, top + Inches(0.15), width, Inches(0.5),
             value, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Consolas')
    add_text(slide, left, top + Inches(0.65), width, Inches(0.3),
             label, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        add_shape(slide, left, top, width or Inches(5), height or Inches(3),
                  fill_color=SURFACE, border_color=BORDER)
        add_text(slide, left, top + (height or Inches(3))//2 - Inches(0.2),
                 width or Inches(5), Inches(0.4),
                 f"[Image: {os.path.basename(path)}]",
                 font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
#  SLIDE GENERATION
# ═══════════════════════════════════════════

prs = new_prs()
BASE = '/home/ec2-user/pai-sim-isaaclab'

# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

add_text(slide, Inches(0), Inches(1.4), SLIDE_W, Inches(0.6),
         "🐾", font_size=40, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(0.8),
         "Physical AI Workshop", font_size=44, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(0.5),
         "$12로 로봇 걷게 만들기", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(0.8),
         "Isaac Lab + PPO로 ANYmal-C 4족 보행 로봇이\n거친 지형에서 걷는 법을 학습하는 전체 과정",
         font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

x_start = Inches(2.8)
for i, (val, lbl, clr) in enumerate([
    ("75분", "훈련 시간", BLUE), ("$12", "총 비용", PURPLE),
    ("4,096", "병렬 환경", PINK), ("1.47억", "timestep", GREEN)
]):
    add_stat_card(slide, x_start + Inches(i * 2.1), Inches(5.0), val, lbl, clr, Inches(1.9))

add_footer(slide, "Physical AI Workshop", "Isaac Lab · PPO · ANYmal-C · AWS g6e.4xlarge")


# ── SLIDE 2: Agenda ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 2)
add_tag(slide, Inches(0.6), Inches(0.4), "📋 AGENDA", PURPLE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "워크샵 구성", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(8), Inches(0.3),
         "7개 Lab + 3개 Appendix · 총 170분", font_size=14, color=TEXT_DIM)

labs = [
    ("Lab 1", "Physical AI 핵심 개념", "10분", BLUE),
    ("Lab 2", "AWS GPU 인프라 구축", "20분", BLUE),
    ("Lab 3", "Isaac Lab Docker 빌드", "30분", PURPLE),
    ("Lab 4", "PPO 강화학습 훈련", "75분", PURPLE),
    ("Lab 5", "학습 결과 분석", "15분", PINK),
    ("Lab 6", "Play 모드 & Export", "10분", PINK),
    ("Lab 7", "정리 & 다음 단계", "10분", GREEN),
]
for i, (lab, desc, time, clr) in enumerate(labs):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(2.0) + Inches(row * 1.2)
    card = add_shape(slide, x, y, Inches(5.8), Inches(1.0), fill_color=SURFACE, border_color=BORDER)
    # Color accent line
    accent = add_shape(slide, x, y, Inches(0.06), Inches(1.0), fill_color=clr)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(4), Inches(0.3),
             f"{lab} · {desc}", font_size=14, color=TEXT, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.55), Inches(4), Inches(0.3),
             time, font_size=11, color=TEXT_MUTED)

# Appendix card
add_shape(slide, Inches(6.8), Inches(5.6), Inches(5.8), Inches(1.0),
          fill_color=SURFACE, border_color=BORDER)
add_text(slide, Inches(7.1), Inches(5.75), Inches(5), Inches(0.3),
         "Appendix A·B·C", font_size=14, color=TEXT_DIM, bold=True)
add_text(slide, Inches(7.1), Inches(6.15), Inches(5), Inches(0.3),
         "트러블슈팅 12선 · 비용 분석 · SW 버전", font_size=11, color=TEXT_MUTED)

add_footer(slide, "Physical AI Workshop", "02 / Agenda")


# ── SLIDE 3: Architecture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 3)
add_tag(slide, Inches(0.6), Inches(0.4), "🏗️ LAB 1", BLUE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "전체 아키텍처", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "Terraform → EC2 GPU → Isaac Lab Docker → PPO 학습 → Policy Export", font_size=14, color=TEXT_DIM)
add_image_safe(slide, f'{BASE}/.gitbook/assets/architecture.png',
               Inches(0.8), Inches(2.0), width=Inches(11.5), height=Inches(4.5))
add_footer(slide, "Physical AI Workshop", "03 / Architecture")


# ── SLIDE 4: Infrastructure ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 4)
add_tag(slide, Inches(0.6), Inches(0.4), "☁️ LAB 2", BLUE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "AWS GPU 인프라", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "Terraform 한 줄로 전체 인프라를 코드로 구축", font_size=14, color=TEXT_DIM)

# Table
rows = [
    ("인스턴스", "g6e.4xlarge (16 vCPU, 64 GB)"),
    ("GPU", "NVIDIA L40S 48 GB VRAM"),
    ("Root EBS", "300 GB gp3 (Docker 이미지)"),
    ("Data EBS", "500 GB gp3 (체크포인트)"),
    ("NVMe", "Instance Store (셰이더 캐시)"),
    ("모니터링", "CloudWatch GPU idle 자동 Stop"),
]
for i, (k, v) in enumerate(rows):
    y = Inches(2.0) + Inches(i * 0.55)
    add_shape(slide, Inches(0.6), y, Inches(5.5), Inches(0.48),
              fill_color=SURFACE if i % 2 == 0 else BG_DARK, border_color=BORDER)
    add_text(slide, Inches(0.8), y + Inches(0.08), Inches(1.8), Inches(0.35),
             k, font_size=12, color=TEXT, bold=True)
    add_text(slide, Inches(2.6), y + Inches(0.08), Inches(3.5), Inches(0.35),
             v, font_size=12, color=TEXT_DIM)

# Code block
code_bg = add_shape(slide, Inches(6.6), Inches(2.0), Inches(6), Inches(4.5),
                    fill_color=RGBColor(0x08, 0x0A, 0x14), border_color=BORDER)
code_lines = [
    ("# 3분 만에 전체 인프라 배포", 11, TEXT_MUTED, False),
    ("terraform init", 13, GREEN, True),
    ("terraform plan", 13, GREEN, True),
    ("terraform apply", 13, GREEN, True),
    ("", 8, TEXT, False),
    ("# SSH 접속", 11, TEXT_MUTED, False),
    ("ssh -i key.pem ubuntu@$(", 12, BLUE, False),
    ("  terraform output -raw public_ip)", 12, BLUE, False),
    ("", 8, TEXT, False),
    ("# 부팅 확인 (15-25분)", 11, TEXT_MUTED, False),
    ("tail -f /var/log/isaac-lab-setup.log", 12, BLUE, False),
    ("# → 'Isaac Lab setup COMPLETE'", 11, TEXT_MUTED, False),
]
add_multiline(slide, Inches(6.9), Inches(2.2), Inches(5.5), Inches(4.2),
              code_lines, font_name='Consolas')

# Cost highlight
cost_shape = add_shape(slide, Inches(0.6), Inches(5.5), Inches(5.5), Inches(0.6),
                       fill_color=RGBColor(0x0A, 0x15, 0x25), border_color=BLUE)
add_text(slide, Inches(0.8), Inches(5.58), Inches(5), Inches(0.4),
         "💰 시간당 ~$3 · 전체 워크샵 ~$12", font_size=14, color=BLUE, bold=True)

add_footer(slide, "Physical AI Workshop", "04 / Infrastructure")


# ── SLIDE 5: Docker Build ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 5)
add_tag(slide, Inches(0.6), Inches(0.4), "🐳 LAB 3", PURPLE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "Isaac Lab Docker 이미지", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "26.8 GB 컨테이너: Isaac Sim 4.5.0 + Isaac Lab v2.1.0 + RSL-RL", font_size=14, color=TEXT_DIM)

# Pitfalls
pitfalls = [
    ("⚠️ NGC에 Isaac Lab 이미지 없음", "→ 소스에서 docker compose 빌드"),
    ("⚠️ 코어 isaaclab 패키지 누락", "→ pip install --no-build-isolation 수동 설치"),
    ("⚠️ 기본 entrypoint 스트리밍 모드", "→ --entrypoint 오버라이드 필수"),
]
for i, (title, desc) in enumerate(pitfalls):
    y = Inches(2.1) + Inches(i * 0.95)
    add_shape(slide, Inches(0.6), y, Inches(6), Inches(0.85),
              fill_color=SURFACE, border_color=RGBColor(0x40, 0x20, 0x10))
    add_text(slide, Inches(0.85), y + Inches(0.1), Inches(5.5), Inches(0.3),
             title, font_size=13, color=ORANGE, bold=True)
    add_text(slide, Inches(0.85), y + Inches(0.45), Inches(5.5), Inches(0.3),
             desc, font_size=12, color=TEXT_DIM)

# Build steps
steps = [
    ("Step 1", "docker compose --profile base build", "~40분"),
    ("Step 2", "pip install -e source/isaaclab + rsl_rl", "~10분"),
    ("Step 3", "docker commit → isaac-lab-ready", "~1분"),
]
for i, (step, cmd, time) in enumerate(steps):
    y = Inches(5.2) + Inches(i * 0.65)
    add_shape(slide, Inches(0.6), y, Inches(6), Inches(0.55),
              fill_color=RGBColor(0x0D, 0x10, 0x28), border_color=RGBColor(0x30, 0x28, 0x50))
    add_text(slide, Inches(0.8), y + Inches(0.1), Inches(0.8), Inches(0.3),
             step, font_size=11, color=GREEN, bold=True, font_name='Consolas')
    add_text(slide, Inches(1.7), y + Inches(0.1), Inches(3.8), Inches(0.3),
             cmd, font_size=11, color=PURPLE, font_name='Consolas')
    add_text(slide, Inches(5.5), y + Inches(0.1), Inches(1), Inches(0.3),
             time, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

add_footer(slide, "Physical AI Workshop", "05 / Docker Build")


# ── SLIDE 6: Training ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 6)
add_tag(slide, Inches(0.6), Inches(0.4), "🏋️ LAB 4", PURPLE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "PPO 강화학습 훈련", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "4,096 병렬 환경 · 1,500 iterations · 75분", font_size=14, color=TEXT_DIM)

phases = [
    ("1", "탐색기", "iter 0-40 · 랜덤 탐색", "-0.5 → -4.9", BLUE),
    ("2", "기초 학습", "iter 40-120 · 보행 패턴 습득", "-4.9 → +5.0", PURPLE),
    ("3", "정교화", "iter 120-300 · 안정적 보행", "+5.0 → +15.0", PINK),
    ("4", "수렴", "iter 300-1500 · 정책 수렴", "+15.0 → +16.3", GREEN),
]
for i, (num, name, desc, result, clr) in enumerate(phases):
    y = Inches(2.0) + Inches(i * 0.85)
    add_shape(slide, Inches(0.6), y, Inches(5.8), Inches(0.75),
              fill_color=RGBColor(clr[0]//10, clr[1]//10, clr[2]//10),
              border_color=RGBColor(clr[0]//5, clr[1]//5, clr[2]//5))
    # Number circle
    num_shape = add_shape(slide, Inches(0.8), y + Inches(0.15), Inches(0.4), Inches(0.4),
                          fill_color=RGBColor(clr[0]//6, clr[1]//6, clr[2]//6))
    add_text(slide, Inches(0.8), y + Inches(0.18), Inches(0.4), Inches(0.35),
             num, font_size=12, color=clr, bold=True, alignment=PP_ALIGN.CENTER, font_name='Consolas')
    add_text(slide, Inches(1.4), y + Inches(0.1), Inches(2), Inches(0.3),
             name, font_size=14, color=clr, bold=True)
    add_text(slide, Inches(1.4), y + Inches(0.4), Inches(2.5), Inches(0.25),
             desc, font_size=10, color=TEXT_MUTED)
    add_text(slide, Inches(4.6), y + Inches(0.2), Inches(1.6), Inches(0.35),
             result, font_size=13, color=clr, bold=True, alignment=PP_ALIGN.RIGHT, font_name='Consolas')

# Highlight
add_shape(slide, Inches(0.6), Inches(5.6), Inches(5.8), Inches(0.55),
          fill_color=RGBColor(0x08, 0x1A, 0x12), border_color=RGBColor(0x15, 0x35, 0x28))
add_text(slide, Inches(0.8), Inches(5.68), Inches(5.5), Inches(0.35),
         "📈 에피소드 길이: 13 → 897 steps (×66배)", font_size=13, color=GREEN, bold=True)

# Dashboard image
add_image_safe(slide, f'{BASE}/.gitbook/assets/dashboard_screenshot.png',
               Inches(6.8), Inches(2.0), width=Inches(5.8), height=Inches(4.2))

add_footer(slide, "Physical AI Workshop", "06 / Training")


# ── SLIDE 7: Results ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 7)
add_tag(slide, Inches(0.6), Inches(0.4), "📊 LAB 5", PINK)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "학습 결과", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "보상 +16.29 · 지형 난이도 5.9/6.25 · VRAM 10/48 GB", font_size=14, color=TEXT_DIM)

# Two dashboard images
add_image_safe(slide, f'{BASE}/.gitbook/assets/image.png',
               Inches(0.6), Inches(2.0), width=Inches(6), height=Inches(3.2))
add_image_safe(slide, f'{BASE}/.gitbook/assets/dashboard_screenshot3.png',
               Inches(6.8), Inches(2.0), width=Inches(6), height=Inches(3.2))

# Stats
stats = [
    ("+16.29", "최종 보상", GREEN), ("897", "에피소드 길이", BLUE),
    ("5.9", "지형 난이도 / 6.25", PURPLE), ("10 GB", "VRAM / 48 GB", PINK),
]
for i, (val, lbl, clr) in enumerate(stats):
    add_stat_card(slide, Inches(0.6 + i * 3.2), Inches(5.6), val, lbl, clr, Inches(2.9))

add_footer(slide, "Physical AI Workshop", "07 / Results")


# ── SLIDE 8: Play Mode ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 8)
add_tag(slide, Inches(0.6), Inches(0.4), "🎬 LAB 6", PINK)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "Play 모드 & Policy Export", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "학습된 정책을 시각화하고, 실제 로봇용으로 변환", font_size=14, color=TEXT_DIM)

# Play image
add_image_safe(slide, f'{BASE}/.gitbook/assets/play30_frame_15s.png',
               Inches(0.6), Inches(2.0), width=Inches(6.2), height=Inches(4.2))
add_text(slide, Inches(0.6), Inches(6.3), Inches(6.2), Inches(0.3),
         "ANYmal-C rough terrain 보행 (30초 비디오)", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Output table
outputs = [
    ("model_1499.pt", "6.6 MB", "학습 재개 / fine-tuning"),
    ("policy.pt", "1.2 MB", "C++ JIT 실시간 추론"),
    ("policy.onnx", "1.1 MB", "TensorRT / Jetson 배포"),
    ("rl-video-*.mp4", "2.8 MB", "시각적 검증"),
]
# Header
add_shape(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.5),
          fill_color=RGBColor(0x10, 0x0D, 0x20), border_color=BORDER)
for j, h in enumerate(["산출물", "크기", "용도"]):
    add_text(slide, Inches(7.4 + j * 1.8), Inches(2.05), Inches(1.8), Inches(0.4),
             h, font_size=10, color=PURPLE, bold=True)
for i, (name, size, use) in enumerate(outputs):
    y = Inches(2.5) + Inches(i * 0.5)
    add_shape(slide, Inches(7.2), y, Inches(5.5), Inches(0.48),
              fill_color=SURFACE if i % 2 == 0 else BG_DARK, border_color=BORDER)
    add_text(slide, Inches(7.4), y + Inches(0.08), Inches(1.8), Inches(0.3),
             name, font_size=11, color=TEXT, bold=True, font_name='Consolas')
    add_text(slide, Inches(9.2), y + Inches(0.08), Inches(1), Inches(0.3),
             size, font_size=11, color=TEXT_DIM)
    add_text(slide, Inches(10.6), y + Inches(0.08), Inches(2), Inches(0.3),
             use, font_size=11, color=TEXT_DIM)

add_footer(slide, "Physical AI Workshop", "08 / Play Mode")


# ── SLIDE 9: Sim-to-Real Pipeline ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 9)
add_tag(slide, Inches(0.6), Inches(0.4), "🤖 SIM-TO-REAL", PURPLE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "Sim-to-Real 파이프라인", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "시뮬레이션에서 학습한 '뇌'를 실제 로봇에 이식", font_size=14, color=TEXT_DIM)
add_image_safe(slide, f'{BASE}/images/sim_to_real_pipeline.png',
               Inches(0.6), Inches(2.0), width=Inches(12), height=Inches(4.7))
add_footer(slide, "Physical AI Workshop", "09 / Sim-to-Real")


# ── SLIDE 10: Deployment ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 10)
add_tag(slide, Inches(0.6), Inches(0.4), "🦿 DEPLOYMENT", PINK)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "ANYmal-C 온보드 배포 구성", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "48차원 관측 → Jetson Orin TensorRT 추론 → 12차원 관절 제어", font_size=14, color=TEXT_DIM)
add_image_safe(slide, f'{BASE}/images/anymal_c_deployment.png',
               Inches(0.6), Inches(2.0), width=Inches(12), height=Inches(4.7))
add_footer(slide, "Physical AI Workshop", "10 / Deployment")


# ── SLIDE 11: Troubleshooting ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 11)
add_tag(slide, Inches(0.6), Inches(0.4), "🔧 APPENDIX", ORANGE)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "실전 트러블슈팅 12선", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "공식 문서에 없는 실전 함정과 해결법", font_size=14, color=TEXT_DIM)

issues = [
    ("1", "dpkg lock 경합", "🟡", ORANGE),
    ("2", "EBS 디바이스 이름 (Nitro)", "🟡", ORANGE),
    ("3", "Instance store 이미 마운트", "🟡", ORANGE),
    ("4", "Terraform templatefile 충돌", "🟡", ORANGE),
    ("5", "user_data 재실행 안 됨", "🟢", GREEN),
    ("6", "Isaac Lab NGC 이미지 없음", "🔴", PINK),
    ("7", "코어 isaaclab 패키지 누락", "🔴", PINK),
    ("8", "Docker entrypoint 스트리밍 모드", "🔴", PINK),
    ("9", "훈련 스크립트 경로 변경", "🟡", ORANGE),
    ("10", "setuptools 빌드 격리", "🔴", PINK),
    ("11", "Volume mount → editable 파괴", "🟡", ORANGE),
    ("12", "셰이더 캐시 첫 실행 지연", "🟢", GREEN),
]
for i, (num, desc, severity, clr) in enumerate(issues):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = Inches(0.6 + col * 6.4)
    y = Inches(2.0) + Inches(row * 0.72)
    add_shape(slide, x, y, Inches(6), Inches(0.62), fill_color=SURFACE, border_color=BORDER)
    add_text(slide, x + Inches(0.15), y + Inches(0.13), Inches(0.35), Inches(0.35),
             num, font_size=11, color=clr, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.6), y + Inches(0.13), Inches(4), Inches(0.35),
             desc, font_size=12, color=TEXT)
    add_text(slide, x + Inches(5.2), y + Inches(0.13), Inches(0.6), Inches(0.35),
             severity, font_size=14, alignment=PP_ALIGN.CENTER)

# Highlight
add_shape(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
          fill_color=RGBColor(0x15, 0x0D, 0x20), border_color=RGBColor(0x30, 0x20, 0x45))
add_text(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.35),
         "💡 🔴 Critical 4개 — 이것을 모르면 훈련 시작조차 할 수 없습니다",
         font_size=13, color=PURPLE, bold=True)

add_footer(slide, "Physical AI Workshop", "11 / Troubleshooting")


# ── SLIDE 12: Cost ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 12)
add_tag(slide, Inches(0.6), Inches(0.4), "💰 COST", GREEN)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "비용 분석", font_size=30, color=TEXT, bold=True)
add_text(slide, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
         "클라우드 GPU로 로봇 AI를 학습하는 실제 비용", font_size=14, color=TEXT_DIM)

cost_items = [
    ("☁️", "이번 실습", "~$12", "g6e.4xlarge On-Demand\nSeoul · 4시간", BLUE),
    ("💚", "비용 최적화", "~$2.50", "g6e.4xlarge Spot\nVirginia 리전", GREEN),
    ("🚀", "대규모", "~$8", "g6e.12xlarge Spot (4×L40S)\nVirginia 리전", PURPLE),
]
for i, (icon, name, cost, detail, clr) in enumerate(cost_items):
    x = Inches(0.8 + i * 4.2)
    add_shape(slide, x, Inches(2.2), Inches(3.8), Inches(2.8),
              fill_color=SURFACE, border_color=RGBColor(clr[0]//4, clr[1]//4, clr[2]//4))
    add_text(slide, x, Inches(2.4), Inches(3.8), Inches(0.5),
             icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.95), Inches(3.8), Inches(0.3),
             name, font_size=16, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.35), Inches(3.8), Inches(0.5),
             cost, font_size=32, color=clr, bold=True, alignment=PP_ALIGN.CENTER, font_name='Consolas')
    add_text(slide, x, Inches(4.0), Inches(3.8), Inches(0.6),
             detail, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Highlights
add_shape(slide, Inches(0.8), Inches(5.4), Inches(5.8), Inches(0.55),
          fill_color=RGBColor(0x0A, 0x15, 0x25), border_color=RGBColor(0x15, 0x28, 0x40))
add_text(slide, Inches(1.0), Inches(5.47), Inches(5.5), Inches(0.4),
         "⏰ GPU 30분 idle → CloudWatch 자동 Stop", font_size=13, color=BLUE, bold=True)

add_shape(slide, Inches(6.8), Inches(5.4), Inches(5.8), Inches(0.55),
          fill_color=RGBColor(0x08, 0x1A, 0x12), border_color=RGBColor(0x12, 0x2E, 0x22))
add_text(slide, Inches(7.0), Inches(5.47), Inches(5.5), Inches(0.4),
         "💾 S3 체크포인트 자동 동기화 (Spot 대비)", font_size=13, color=GREEN, bold=True)

add_footer(slide, "Physical AI Workshop", "12 / Cost")


# ── SLIDE 13: Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_slide_num(slide, 13)
add_tag(slide, Inches(0.6), Inches(0.4), "✅ SUMMARY", GREEN)
add_text(slide, Inches(0.6), Inches(0.9), Inches(8), Inches(0.5),
         "이번 워크샵에서 달성한 것", font_size=30, color=TEXT, bold=True)

achievements = [
    "✅ Terraform으로 AWS GPU 인프라를 코드로 구축",
    "✅ Isaac Sim + Isaac Lab Docker 이미지 빌드 (26.8 GB)",
    "✅ 4,096개 병렬 환경에서 PPO 강화학습 실행",
    "✅ 75분 만에 1.47억 timestep 학습 완료",
    "✅ ANYmal-C가 rough terrain에서 안정적으로 보행",
    "✅ 학습된 정책을 JIT/ONNX로 export (sim-to-real 준비)",
    "✅ 총 비용: ~$12 (₩16,000)",
]
for i, text in enumerate(achievements):
    add_text(slide, Inches(0.8), Inches(1.7 + i * 0.52), Inches(11), Inches(0.45),
             text, font_size=16, color=TEXT)

# Pipeline
pipeline = [
    ("🏗️", "인프라"), ("🐳", "Docker"), ("🏋️", "훈련"),
    ("📊", "분석"), ("🎬", "Play"), ("🤖", "Sim-to-Real"),
]
for i, (icon, name) in enumerate(pipeline):
    x = Inches(0.8 + i * 2.0)
    y = Inches(5.8)
    add_shape(slide, x, y, Inches(1.5), Inches(0.9), fill_color=SURFACE, border_color=BORDER)
    add_text(slide, x, y + Inches(0.08), Inches(1.5), Inches(0.4),
             icon, font_size=18, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.48), Inches(1.5), Inches(0.3),
             name, font_size=11, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_text(slide, x + Inches(1.55), y + Inches(0.25), Inches(0.4), Inches(0.35),
                 "→", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_footer(slide, "Physical AI Workshop", "13 / Summary")


# ── SLIDE 14: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.6),
         "🐾", font_size=48, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.3), SLIDE_W, Inches(0.8),
         "감사합니다", font_size=44, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.5),
         "$12로 4족 보행 로봇이 걷는 법을 배웠습니다", font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

for i, (val, lbl) in enumerate([("Isaac Lab", "시뮬레이션"), ("PPO", "알고리즘"), ("ANYmal-C", "로봇")]):
    x = Inches(3.5 + i * 2.3)
    add_stat_card(slide, x, Inches(4.3), val, lbl,
                  [BLUE, PURPLE, PINK][i], Inches(2.0))

add_text(slide, Inches(0), Inches(6.0), SLIDE_W, Inches(0.3),
         "GitHub · comeddy/pai-sim-isaaclab", font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════
output_path = f'{BASE}/Physical_AI_Workshop.pptx'
prs.save(output_path)
print(f'✅ Saved: {output_path}')
print(f'   Slides: {len(prs.slides)}')
print(f'   Size: {os.path.getsize(output_path) / 1024:.0f} KB')
